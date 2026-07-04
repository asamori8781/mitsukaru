"""ファイル形式ごとのテキスト抽出(Phase 1: 全文インデックス用)。

ここで抽出したテキストはローカルのDB(file_content/file_content_fts)にのみ
保存され、外部には送信しない(ai_client.pyとは完全に独立している)。
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

EXTRACTOR_VERSION = 1

# indexer.TEXT_LIKE_EXTENSIONSと同じ「そのまま読める」テキスト系拡張子。
# ここでは文字コード判定に少し幅を持たせて直接読み込む。
PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".log", ".ini",
    ".yaml", ".yml", ".py", ".js", ".ts", ".css", ".java", ".c", ".cpp", ".h",
    ".hpp", ".cs", ".go", ".rb", ".php", ".sh", ".bat", ".ps1", ".sql", ".rtf",
}

MAX_CHARS = 300_000  # 1ファイルあたりの抽出テキスト上限(異常に大きい文書の暴走防止)


class ExtractionError(Exception):
    pass


def _truncate(text: str) -> str:
    return text[:MAX_CHARS]


def _extract_plain_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp932", "shift_jis"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    import pypdf

    reader = pypdf.PdfReader(str(path))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as e:
            raise ExtractionError(f"暗号化されたPDFを開けませんでした: {e}") from e
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        parts.append(str(cell))
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    parts.append("".join(run.text for run in paragraph.runs))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def is_extractable(ext: str) -> bool:
    return ext in PLAIN_TEXT_EXTENSIONS or ext in _EXTRACTORS


def extractable_extensions() -> set[str]:
    return PLAIN_TEXT_EXTENSIONS | set(_EXTRACTORS.keys())


def extract_text(path: str, ext: str) -> Optional[str]:
    """テキストを抽出する。対応外の拡張子はNoneを返す。抽出失敗はExtractionErrorを送出する。"""
    p = Path(path)
    try:
        if ext in PLAIN_TEXT_EXTENSIONS:
            text = _extract_plain_text(p)
        elif ext in _EXTRACTORS:
            text = _EXTRACTORS[ext](p)
        else:
            return None
    except ExtractionError:
        raise
    except Exception as e:
        raise ExtractionError(f"{ext}ファイルのテキスト抽出に失敗しました: {e}") from e
    return _truncate(text)
